"""Text extraction for multiple document formats.

Every extractor returns the same shape the ingestion pipeline expects:
    list[tuple[int, str]]   # [(section_number, text), ...]

For PDFs "section" is the page; for PPTX it's the slide; for DOCX it's a
running block index; for XLSX/CSV it's the sheet. That keeps citations
meaningful ("page 3", "slide 5") across formats without changing the
downstream chunk/embed/store code.
"""
import csv
import io

import fitz  # PyMuPDF (PDF)


class ExtractionError(Exception):
    """Raised when a file cannot be read."""


# Supported extensions -> friendly label
SUPPORTED = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".txt": "Text",
    ".md": "Markdown",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".csv": "CSV",
}


def extension_of(filename: str) -> str:
    name = filename.lower()
    dot = name.rfind(".")
    return name[dot:] if dot != -1 else ""


def is_supported(filename: str) -> bool:
    return extension_of(filename) in SUPPORTED


# ---------------- per-format extractors ----------------

def _extract_pdf(data: bytes) -> list[tuple[int, str]]:
    pages = []
    with fitz.open(stream=io.BytesIO(data), filetype="pdf") as pdf:
        for i, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append((i, text))
    return pages


def _extract_docx(data: bytes) -> list[tuple[int, str]]:
    from docx import Document
    doc = Document(io.BytesIO(data))
    # Group paragraphs into pseudo-pages of ~40 non-empty paragraphs so
    # citations point at a rough location rather than one giant block.
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # include table text too
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    if not paras:
        return []
    sections, chunk, page = [], [], 1
    for i, p in enumerate(paras, start=1):
        chunk.append(p)
        if i % 40 == 0:
            sections.append((page, "\n".join(chunk)))
            chunk, page = [], page + 1
    if chunk:
        sections.append((page, "\n".join(chunk)))
    return sections


def _extract_text(data: bytes) -> list[tuple[int, str]]:
    text = data.decode("utf-8", errors="replace").strip()
    return [(1, text)] if text else []


def _extract_pptx(data: bytes) -> list[tuple[int, str]]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            slides.append((i, "\n".join(parts)))
    return slides


def _extract_xlsx(data: bytes) -> list[tuple[int, str]]:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for idx, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            header = f"Sheet: {ws.title}"
            sheets.append((idx, header + "\n" + "\n".join(rows)))
    wb.close()
    return sheets


def _extract_csv(data: bytes) -> list[tuple[int, str]]:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(cell.strip() for cell in row) for row in reader if any(row)]
    return [(1, "\n".join(rows))] if rows else []


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".txt": _extract_text,
    ".md": _extract_text,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
}


def extract_sections(filename: str, data: bytes) -> list[tuple[int, str]]:
    """Extract [(section_number, text)] for any supported format."""
    ext = extension_of(filename)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ExtractionError(f"Unsupported file type: {ext or 'unknown'}")
    try:
        sections = extractor(data)
    except ExtractionError:
        raise
    except Exception as e:
        raise ExtractionError(f"Could not read {SUPPORTED.get(ext, ext)} file: {e}") from e

    if not sections:
        raise ExtractionError(
            "No extractable text found. The file may be empty, image-only, "
            "or scanned."
        )
    return sections


def count_sections(filename: str, data: bytes) -> int:
    """Section/page count for plan validation, without full extraction where
    possible."""
    ext = extension_of(filename)
    if ext == ".pdf":
        try:
            with fitz.open(stream=io.BytesIO(data), filetype="pdf") as pdf:
                return pdf.page_count
        except Exception as e:
            raise ExtractionError(f"Could not read PDF: {e}") from e
    # For other formats, fall back to counting extracted sections.
    return len(extract_sections(filename, data))
